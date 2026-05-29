"""
Generate a 3-slide PPTX from the Shawan digital city project content.
Design: professional government-report style with dark blue + red accents.
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu, Cm
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import copy

# === Color palette ===
INK      = RGBColor(0x1A, 0x1A, 0x1A)
BLUE     = RGBColor(0x0D, 0x3B, 0x66)
RED      = RGBColor(0xB7, 0x1C, 0x1C)
GOLD     = RGBColor(0x8B, 0x69, 0x14)
WHITE    = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_BG = RGBColor(0xF7, 0xF6, 0xF2)
GRAY     = RGBColor(0x88, 0x88, 0x88)
DARK_GRAY = RGBColor(0x55, 0x55, 0x55)
BORDER   = RGBColor(0xDD, 0xD9, 0xCF)
TABLE_HDR = BLUE
TABLE_STRIPE = RGBColor(0xFA, 0xFA, 0xF8)
SUM_BG   = RGBColor(0xF5, 0xF2, 0xEA)
NOTE_BG  = RGBColor(0xFE, 0xFC, 0xF5)
TAG_RED_BG = RGBColor(0xFE, 0xF5, 0xF5)
TAG_BLUE_BG = RGBColor(0xE8, 0xF0, 0xF8)
TAG_GOLD_BG = RGBColor(0xFD, 0xF3, 0xD1)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

prs = Presentation()
prs.slide_width  = SLIDE_W
prs.slide_height = SLIDE_H

# Use blank layout
blank_layout = prs.slide_layouts[6]  # blank


# === Helper functions ===

def add_rect(slide, left, top, width, height, fill_color=None, border_color=None):
    """Add a rectangle shape."""
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.line.fill.background()
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(0.5)
    return shape

def add_textbox(slide, left, top, width, height, text="", font_size=Pt(12),
                color=INK, bold=False, align=PP_ALIGN.LEFT, font_name='Microsoft YaHei'):
    """Add a text box with a single paragraph."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = font_size
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = align
    return tf

def add_rich_textbox(slide, left, top, width, height):
    """Add a text box and return its text_frame for manual paragraph building."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    return tf

def add_para(tf, text, size=Pt(12), color=INK, bold=False, align=PP_ALIGN.LEFT,
             space_before=Pt(0), space_after=Pt(0), font_name='Microsoft YaHei'):
    """Add a paragraph to a text frame."""
    p = tf.add_paragraph()
    p.text = text
    p.font.size = size
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = align
    p.space_before = space_before
    p.space_after = space_after
    return p

def first_para(tf, text, size=Pt(12), color=INK, bold=False, align=PP_ALIGN.LEFT,
               font_name='Microsoft YaHei'):
    """Set text on the first (default) paragraph."""
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = size
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = align
    return p

def add_run_to_para(p, text, size=Pt(12), color=INK, bold=False, font_name='Microsoft YaHei'):
    """Add a run to an existing paragraph."""
    run = p.add_run()
    run.text = text
    run.font.size = size
    run.font.color.rgb = color
    run.font.bold = bold
    run.font.name = font_name
    return run

def add_slide_number(slide, num, total=3):
    """Add page number badge top-right."""
    shape = add_rect(slide, Inches(11.6), Inches(0.28), Inches(1.2), Inches(0.32), fill_color=BLUE)
    shape.line.fill.background()
    tf = shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = f"{num:02d} / {total:02d}"
    p.font.size = Pt(9)
    p.font.color.rgb = WHITE
    p.font.bold = False
    p.font.name = 'Microsoft YaHei'
    p.alignment = PP_ALIGN.CENTER

def add_top_bar(slide):
    """Add the colored gradient bar at the very top of the slide."""
    # We'll simulate gradient with 3 thin rectangles
    add_rect(slide, Inches(0), Inches(0), Inches(4.44), Inches(0.06), fill_color=RED)
    add_rect(slide, Inches(4.44), Inches(0), Inches(4.44), Inches(0.06), fill_color=BLUE)
    add_rect(slide, Inches(8.88), Inches(0), Inches(4.453), Inches(0.06), fill_color=GOLD)

def add_slide_header(slide, title_text, subtitle_lines):
    """Add the slide header with title + subtitle, plus bottom border line."""
    # Title
    add_textbox(slide, Inches(0.7), Inches(0.4), Inches(8), Inches(0.5),
                title_text, font_size=Pt(24), color=BLUE, bold=True)
    # Subtitle
    add_textbox(slide, Inches(8.5), Inches(0.4), Inches(4), Inches(0.5),
                subtitle_lines, font_size=Pt(10), color=GRAY, align=PP_ALIGN.RIGHT)
    # Bottom border line
    add_rect(slide, Inches(0.7), Inches(1.0), Inches(11.9), Inches(0.02), fill_color=BLUE)

def add_footer(slide):
    """Add footer bar at bottom."""
    add_rect(slide, Inches(0.7), Inches(7.05), Inches(11.9), Inches(0.01), fill_color=RGBColor(0xEC, 0xE9, 0xE0))
    add_textbox(slide, Inches(0.7), Inches(7.08), Inches(5), Inches(0.3),
                "乐山市沙湾区", font_size=Pt(8), color=RGBColor(0xAA, 0xAA, 0xAA))
    add_textbox(slide, Inches(8.5), Inches(7.08), Inches(4), Inches(0.3),
                "2026 年 5 月", font_size=Pt(8), color=RGBColor(0xAA, 0xAA, 0xAA), align=PP_ALIGN.RIGHT)

def add_section_label(slide, left, top, text):
    """Add a small red section label."""
    add_textbox(slide, left, top, Inches(4), Inches(0.22),
                text, font_size=Pt(9), color=RED, bold=True)

def make_table(slide, left, top, col_widths, headers, rows, font_size=Pt(10)):
    """Create a styled table. Returns the table shape."""
    n_rows = len(rows) + 1  # +1 for header
    n_cols = len(headers)
    total_w = sum(col_widths)

    table_shape = slide.shapes.add_table(n_rows, n_cols, left, top, total_w, Inches(0.32 * n_rows))
    table = table_shape.table

    for ci, cw in enumerate(col_widths):
        table.columns[ci].width = cw

    # Header row
    for ci, hdr in enumerate(headers):
        cell = table.cell(0, ci)
        cell.text = hdr
        cell.fill.solid()
        cell.fill.fore_color.rgb = TABLE_HDR
        for p in cell.text_frame.paragraphs:
            p.font.size = Pt(9)
            p.font.color.rgb = WHITE
            p.font.bold = True
            p.font.name = 'Microsoft YaHei'
            p.alignment = PP_ALIGN.CENTER if ci > 0 else PP_ALIGN.LEFT

    # Data rows
    for ri, row in enumerate(rows):
        for ci, val in enumerate(row):
            cell = table.cell(ri + 1, ci)
            cell.text = str(val) if val is not None else "—"
            if ri % 2 == 1:
                cell.fill.solid()
                cell.fill.fore_color.rgb = TABLE_STRIPE
            else:
                cell.fill.solid()
                cell.fill.fore_color.rgb = WHITE
            for p in cell.text_frame.paragraphs:
                p.font.size = font_size
                p.font.color.rgb = INK
                p.font.name = 'Microsoft YaHei'
                p.alignment = PP_ALIGN.CENTER if ci > 0 else PP_ALIGN.LEFT

    return table_shape


# ============================================================
# SLIDE 1: 为什么要建？
# ============================================================
s1 = prs.slides.add_slide(blank_layout)
add_top_bar(s1)
add_slide_number(s1, 1)
add_slide_header(s1, "为什么要建？",
                 "乐山市沙湾区\n数字城市市政公共服务信息化建设项目")

# --- Policy cards ---
add_section_label(s1, Inches(0.7), Inches(1.2), "政策要求")

policy_data = [
    ("国家层",
     "数字经济规模 63.2 万亿，GDP 贡献率 66.45%；数据要素×三年行动覆盖 12 个重点领域；数字中国\"2522\"框架全面部署。"),
    ("四川省",
     "2027 年建成全省数字政府总体架构；智慧城市按\"1+1+N\"架构推进，市县两级全域数字化转型。"),
    ("乐山市",
     "智慧城市考核覆盖 10 大类 38 项指标，涵盖惠民服务、精准治理、生态宜居、数据安全等，对标全国文明城市标准。"),
]

card_w = Inches(3.75)
card_gap = Inches(0.18)
card_left_start = Inches(0.7)
card_top = Inches(1.5)

for i, (level, desc) in enumerate(policy_data):
    x = card_left_start + i * (card_w + card_gap)
    # Card background
    add_rect(s1, x, card_top, card_w, Inches(1.05), fill_color=LIGHT_BG)
    # Left accent bar
    add_rect(s1, x, card_top, Inches(0.04), Inches(1.05), fill_color=BLUE)
    # Level label
    add_textbox(s1, x + Inches(0.2), card_top + Inches(0.08), Inches(3), Inches(0.25),
                level, font_size=Pt(11), color=RED, bold=True)
    # Description
    add_textbox(s1, x + Inches(0.2), card_top + Inches(0.35), Inches(3.3), Inches(0.7),
                desc, font_size=Pt(9.5), color=INK)

# --- Pain points ---
add_section_label(s1, Inches(0.7), Inches(2.8), "现状痛点")
add_textbox(s1, Inches(0.7), Inches(3.0), Inches(6), Inches(0.35),
            "四大问题倒逼数字化转型", font_size=Pt(16), color=INK, bold=True)

pain_data = [
    ("1", "治理手段单一", "传统模式效率低\n人为差错率高", "健全全域感知", RED, TAG_RED_BG),
    ("2", "业务应用不足", "专项领域场景缺失\n数字化覆盖低", "统筹融合应用", BLUE, TAG_BLUE_BG),
    ("3", "数据底座薄弱", "系统\"烟囱\"林立\n数据汇聚复用差", "统一数据底座", GOLD, TAG_GOLD_BG),
    ("4", "跨部门协同难", "信息流转不畅\n指挥调度分散", "统一运营指挥", RED, TAG_RED_BG),
]

pain_w = Inches(2.85)
pain_gap = Inches(0.15)
pain_top = Inches(3.45)
pain_left = Inches(0.7)

for i, (num, name, detail, fix, accent_color, tag_bg) in enumerate(pain_data):
    x = pain_left + i * (pain_w + pain_gap)
    # Card background
    add_rect(s1, x, pain_top, pain_w, Inches(2.8), fill_color=WHITE, border_color=BORDER)

    # Number circle
    circle = s1.shapes.add_shape(MSO_SHAPE.OVAL, x + Inches(1.1), pain_top + Inches(0.15),
                                  Inches(0.5), Inches(0.5))
    circle.fill.solid()
    circle.fill.fore_color.rgb = accent_color
    circle.line.fill.background()
    tf = circle.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.text = num
    p.font.size = Pt(16)
    p.font.color.rgb = WHITE
    p.font.bold = True
    p.font.name = 'Microsoft YaHei'
    p.alignment = PP_ALIGN.CENTER

    # Pain name
    add_textbox(s1, x + Inches(0.15), pain_top + Inches(0.8), Inches(2.55), Inches(0.3),
                name, font_size=Pt(13), color=INK, bold=True, align=PP_ALIGN.CENTER)
    # Pain detail
    add_textbox(s1, x + Inches(0.15), pain_top + Inches(1.2), Inches(2.55), Inches(0.55),
                detail, font_size=Pt(9.5), color=DARK_GRAY, align=PP_ALIGN.CENTER)
    # Fix tag
    tag_shape = add_rect(s1, x + Inches(0.45), pain_top + Inches(1.95), Inches(1.95), Inches(0.28),
                         fill_color=tag_bg)
    tf2 = tag_shape.text_frame
    tf2.word_wrap = True
    p2 = tf2.paragraphs[0]
    p2.text = fix
    p2.font.size = Pt(9)
    p2.font.color.rgb = accent_color
    p2.font.bold = False
    p2.font.name = 'Microsoft YaHei'
    p2.alignment = PP_ALIGN.CENTER

add_footer(s1)


# ============================================================
# SLIDE 2: 怎么建？
# ============================================================
s2 = prs.slides.add_slide(blank_layout)
add_top_bar(s2)
add_slide_number(s2, 2)
add_slide_header(s2, "怎么建？",
                 "乐山市沙湾区\n数字城市市政公共服务信息化建设项目")

# --- Concept bar ---
add_section_label(s2, Inches(0.7), Inches(1.2), "总体思路")
concept_shape = add_rect(s2, Inches(0.7), Inches(1.45), Inches(11.9), Inches(0.48),
                         fill_color=LIGHT_BG, border_color=BORDER)
tf = concept_shape.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.alignment = PP_ALIGN.CENTER
add_run_to_para(p, "紧扣\"", size=Pt(13), color=INK)
add_run_to_para(p, "工业强区、农旅兴区", size=Pt(13), color=INK, bold=True)
add_run_to_para(p, "\"发展主线  →  ", size=Pt(13), color=INK)
add_run_to_para(p, "以产兴业", size=Pt(13), color=RED, bold=True)
add_run_to_para(p, " 强经济 · ", size=Pt(13), color=INK)
add_run_to_para(p, "以城赋能", size=Pt(13), color=RED, bold=True)
add_run_to_para(p, " 优治理 · ", size=Pt(13), color=INK)
add_run_to_para(p, "以人为本", size=Pt(13), color=RED, bold=True)
add_run_to_para(p, " 惠民生", size=Pt(13), color=INK)

# --- Architecture ---
add_section_label(s2, Inches(0.7), Inches(2.15), "总体架构")
add_textbox(s2, Inches(0.7), Inches(2.35), Inches(8), Inches(0.3),
            "\"三层架构\"支撑全域数字化转型", font_size=Pt(16), color=INK, bold=True)

arch_top = Inches(2.8)
arch_h = Inches(3.2)

# Column widths
col_w = [Inches(3.2), Inches(5.1), Inches(3.6)]
col_gap = Inches(0.15)
col_heads = ["主题应用层 — 三线并进", "数字大脑底座 — 五大核心能力", "基础设施层 — 感知 + 承载"]
col_colors = [RED, BLUE, RGBColor(0x3A, 0x60, 0x80)]

for ci in range(3):
    x = Inches(0.7) + ci * (col_w[ci] + col_gap)
    # Column container
    add_rect(s2, x, arch_top, col_w[ci], arch_h, fill_color=WHITE, border_color=BORDER)
    # Column header
    add_rect(s2, x, arch_top, col_w[ci], Inches(0.35), fill_color=col_colors[ci])
    add_textbox(s2, x + Inches(0.1), arch_top + Inches(0.03), col_w[ci] - Inches(0.2), Inches(0.3),
                col_heads[ci], font_size=Pt(10), color=WHITE, bold=True, align=PP_ALIGN.CENTER)

# Column 1: Application layer
app_items = [
    ("产 · 数字经济", "智慧文旅 / 智慧招商 / 经济运行态势监测"),
    ("城 · 数字政府", "智慧交通 / 智慧应急 / 城市治理全域管控 / 应急指挥速处联防"),
    ("人 · 数字社会", "智慧医疗 / 智慧教育 / 民生服务智慧保障"),
]
y0 = arch_top + Inches(0.5)
for j, (title, subs) in enumerate(app_items):
    y = y0 + j * Inches(0.85)
    item_shape = add_rect(s2, Inches(0.8), y, Inches(2.9), Inches(0.75),
                          fill_color=LIGHT_BG, border_color=BLUE)
    tf2 = item_shape.text_frame
    tf2.word_wrap = True
    first_para(tf2, title, size=Pt(11), color=INK, bold=True)
    add_para(tf2, subs, size=Pt(9), color=DARK_GRAY)

# Column 2: Base layer
base_items = [
    ("智能感知", "视频融合 · 物联网平台"),
    ("数据集成交换", "数据治理 · 数据共享 · 统一数据出口"),
    ("时空信息服务", "CIM 模型引擎 · 空天地数据管理"),
    ("业务支撑平台", "API 管理 · 可视化 · 权限管理 · 能力开放"),
    ("智能服务中心", "算法平台 · AI 服务 · AI 运营"),
]
x2 = Inches(0.7) + col_w[0] + col_gap
for j, (title, subs) in enumerate(base_items):
    y = y0 + j * Inches(0.52)
    item_shape = add_rect(s2, x2 + Inches(0.1), y, Inches(4.9), Inches(0.44),
                          fill_color=LIGHT_BG, border_color=BLUE)
    tf2 = item_shape.text_frame
    tf2.word_wrap = True
    first_para(tf2, "", size=Pt(10))
    # Can't mix bold in same paragraph easily, just set first text
    p = tf2.paragraphs[0]
    p.clear()
    add_run_to_para(p, title, size=Pt(10), color=INK, bold=True)
    add_run_to_para(p, f"  — {subs}", size=Pt(9), color=DARK_GRAY)

# Column 3: Infrastructure
infra_items = [
    ("全域感知源", "教育设施 · 交通视频图像感知源 · 智慧文旅机器人 ·\n智慧管网感知源 · 环保感知源"),
    ("运营管理中心", "指挥大厅 · 计算资源 · 基础网络 · 机房配套"),
]
x3 = Inches(0.7) + col_w[0] + col_gap + col_w[1] + col_gap
for j, (title, subs) in enumerate(infra_items):
    y = y0 + j * Inches(1.2)
    item_shape = add_rect(s2, x3 + Inches(0.1), y, Inches(3.4), Inches(1.05),
                          fill_color=LIGHT_BG, border_color=BLUE)
    tf2 = item_shape.text_frame
    tf2.word_wrap = True
    first_para(tf2, title, size=Pt(11), color=INK, bold=True)
    add_para(tf2, subs, size=Pt(9), color=DARK_GRAY)

# --- Goal strip ---
goal_y = Inches(6.2)
goal_shape = add_rect(s2, Inches(0.7), goal_y, Inches(11.9), Inches(0.6), fill_color=BLUE)
tf = goal_shape.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "最终目标：沙湾数字城市运营一网统管"
p.font.size = Pt(16)
p.font.color.rgb = WHITE
p.font.bold = True
p.font.name = 'Microsoft YaHei'
p.alignment = PP_ALIGN.CENTER
add_para(tf, "\"一屏观全域、一网管全城、一体指挥调度\"", size=Pt(9), color=RGBColor(0xCC, 0xCC, 0xCC),
         align=PP_ALIGN.CENTER)

add_footer(s2)


# ============================================================
# SLIDE 3: 花多少钱？值不值？
# ============================================================
s3 = prs.slides.add_slide(blank_layout)
add_top_bar(s3)
add_slide_number(s3, 3)
add_slide_header(s3, "花多少钱？值不值？",
                 "乐山市沙湾区\n数字城市市政公共服务信息化建设项目")

# --- Left: Investment table ---
add_section_label(s3, Inches(0.7), Inches(1.2), "投资概览")
add_textbox(s3, Inches(0.7), Inches(1.42), Inches(7), Inches(0.3),
            "总投资 1.39 亿元，已设计 7,210 万元", font_size=Pt(16), color=INK, bold=True)

invest_headers = ["建设内容", "原初设（万元）", "已设计（万元）"]
invest_rows = [
    ["数字城市业务支撑平台", "1,656", "—"],
    ["数字沙湾运营平台",     "308",   "—"],
    ["智慧教育",             "1,846", "2,508"],
    ["智慧医疗",             "2,690", "—"],
    ["应急指挥",             "—",     "150"],
    ["云网租赁费",           "3,431", "2,251"],
    ["系统集成费",           "496",   "185"],
]
invest_col_w = [Inches(3.6), Inches(1.6), Inches(1.8)]
tbl_shape = make_table(s3, Inches(0.7), Inches(1.85), invest_col_w, invest_headers, invest_rows, font_size=Pt(10))

# Add sum row manually by modifying the table - actually, let's add a separate row shape below
# Or just add a text line. Actually let me add the sum as a separate visual element.
sum_y = Inches(1.85) + Inches(0.32 * 8)  # 8 rows total (1 header + 7 data)
# Actually let's add the sum row to the table
# The table already has the header + rows. Let me add the sum below.
sum_shape = add_rect(s3, Inches(0.7), sum_y, sum(invest_col_w), Inches(0.32), fill_color=SUM_BG)
sum_shape.line.color.rgb = BLUE
sum_shape.line.width = Pt(1.5)

# Add "合计" text
add_textbox(s3, Inches(0.8), sum_y + Inches(0.02), Inches(3.4), Inches(0.28),
            "合计", font_size=Pt(10), color=INK, bold=True)
add_textbox(s3, Inches(4.4), sum_y + Inches(0.02), Inches(1.4), Inches(0.28),
            "13,851", font_size=Pt(10), color=INK, bold=True, align=PP_ALIGN.CENTER)
add_textbox(s3, Inches(5.95), sum_y + Inches(0.02), Inches(1.6), Inches(0.28),
            "7,210", font_size=Pt(10), color=INK, bold=True, align=PP_ALIGN.CENTER)

# Note
note_y = sum_y + Inches(0.45)
note_shape = add_rect(s3, Inches(0.7), note_y, Inches(6.8), Inches(0.45), fill_color=NOTE_BG,
                      border_color=RGBColor(0xE5, 0xD9, 0xA0))
tf = note_shape.text_frame
tf.word_wrap = True
first_para(tf, "", size=Pt(9))
p = tf.paragraphs[0]
p.clear()
add_run_to_para(p, "变动说明：", size=Pt(9), color=GOLD, bold=True)
add_run_to_para(p, "智慧交通较初设缩减 ", size=Pt(9), color=GOLD)
add_run_to_para(p, "36.8%", size=Pt(9), color=GOLD, bold=True)
add_run_to_para(p, "；智慧教育增加 ", size=Pt(9), color=GOLD)
add_run_to_para(p, "35.9%", size=Pt(9), color=GOLD, bold=True)
add_run_to_para(p, "（询价中）。部分超概项需进一步询价、核减规模。", size=Pt(9), color=GOLD)

# --- Right: Outcomes ---
out_x = Inches(8.0)
add_section_label(s3, out_x, Inches(1.2), "预期成效")
add_textbox(s3, out_x, Inches(1.42), Inches(4.8), Inches(0.3),
            "三维价值回报", font_size=Pt(16), color=INK, bold=True)

outcomes = [
    ("1", "治理效能升级",
     "\"大综治\"模式落地，全域感知 + 数据驱动决策，城市事件从被动响应转向主动预警，一网统管提效。"),
    ("2", "产业数字赋能",
     "贯通工业、农业、文旅三大产业链，以数据要素流通带动招商提效、经济态势可测、产业竞争力提升。"),
    ("3", "民生普惠提质",
     "智慧医疗、智慧教育覆盖全域，数字化服务均等化，切实增强居民获得感与满意度。"),
]

for j, (num, title, desc) in enumerate(outcomes):
    y = Inches(2.0) + j * Inches(1.5)
    # Number circle
    circle = s3.shapes.add_shape(MSO_SHAPE.OVAL, out_x, y, Inches(0.4), Inches(0.4))
    circle.fill.solid()
    circle.fill.fore_color.rgb = RED
    circle.line.fill.background()
    tf = circle.text_frame
    p = tf.paragraphs[0]
    p.text = num
    p.font.size = Pt(14)
    p.font.color.rgb = WHITE
    p.font.bold = True
    p.font.name = 'Microsoft YaHei'
    p.alignment = PP_ALIGN.CENTER

    # Title + desc
    add_textbox(s3, out_x + Inches(0.55), y - Inches(0.02), Inches(4), Inches(0.25),
                title, font_size=Pt(13), color=INK, bold=True)
    add_textbox(s3, out_x + Inches(0.55), y + Inches(0.28), Inches(4), Inches(0.8),
                desc, font_size=Pt(10), color=DARK_GRAY)

    # Separator line (except last)
    if j < 2:
        add_rect(s3, out_x + Inches(0.55), y + Inches(1.2), Inches(4.3), Inches(0.005),
                 fill_color=RGBColor(0xEE, 0xEE, 0xEE))

add_footer(s3)


# === Save ===
output_path = r"E:\claude\output\沙湾数字城市汇报_3页.pptx"
prs.save(output_path)
print(f"PPTX saved to: {output_path}")
