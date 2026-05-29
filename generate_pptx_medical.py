"""
Generate a 3-slide PPTX for Shawan Smart Medical survey & planning.
Lesson learned: align Y-coordinates across multi-column items using a unified row grid.
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Cm, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# === Colors ===
INK       = RGBColor(0x1A, 0x1A, 0x1A)
BLUE      = RGBColor(0x0D, 0x3B, 0x66)
RED       = RGBColor(0xB7, 0x1C, 0x1C)
GOLD      = RGBColor(0x8B, 0x69, 0x14)
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_BG  = RGBColor(0xF7, 0xF6, 0xF2)
GRAY      = RGBColor(0x88, 0x88, 0x88)
DARK_GRAY = RGBColor(0x55, 0x55, 0x55)
BORDER    = RGBColor(0xDD, 0xD9, 0xCF)
TAG_RED   = RGBColor(0xFE, 0xF5, 0xF5)
TAG_BLUE  = RGBColor(0xE8, 0xF0, 0xF8)
TAG_GOLD  = RGBColor(0xFE, 0xF9, 0xE7)
SUM_BG    = RGBColor(0xF5, 0xF2, 0xEA)
NOTE_BG   = RGBColor(0xFE, 0xFC, 0xF5)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

prs = Presentation()
prs.slide_width  = SLIDE_W
prs.slide_height = SLIDE_H
blank = prs.slide_layouts[6]

# ===================== HELPERS =====================

def rect(slide, l, t, w, h, fill=None, border=None, border_w=None):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, l, t, w, h)
    s.line.fill.background()
    if fill:
        s.fill.solid(); s.fill.fore_color.rgb = fill
    if border:
        s.line.color.rgb = border
        s.line.width = border_w or Pt(0.5)
    return s

def oval(slide, l, t, w, h, fill=None):
    s = slide.shapes.add_shape(MSO_SHAPE.OVAL, l, t, w, h)
    s.line.fill.background()
    if fill:
        s.fill.solid(); s.fill.fore_color.rgb = fill
    return s

def tb(slide, l, t, w, h, text="", size=Pt(12), color=INK, bold=False, align=PP_ALIGN.LEFT, font='Microsoft YaHei'):
    """Simple single-text textbox. Returns text_frame."""
    txBox = slide.shapes.add_textbox(l, t, w, h)
    tf = txBox.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.text = text
    p.font.size = size; p.font.color.rgb = color; p.font.bold = bold
    p.font.name = font; p.alignment = align
    return tf

def rich_tb(slide, l, t, w, h):
    txBox = slide.shapes.add_textbox(l, t, w, h)
    tf = txBox.text_frame; tf.word_wrap = True
    return tf

def run(p, text, size=Pt(12), color=INK, bold=False, font='Microsoft YaHei'):
    r = p.add_run(); r.text = text
    r.font.size = size; r.font.color.rgb = color; r.font.bold = bold
    r.font.name = font
    return r

def add_para(tf, text, size=Pt(12), color=INK, bold=False, align=PP_ALIGN.LEFT, sb=Pt(0), sa=Pt(0), font='Microsoft YaHei'):
    p = tf.add_paragraph(); p.text = text
    p.font.size = size; p.font.color.rgb = color; p.font.bold = bold
    p.font.name = font; p.alignment = align
    p.space_before = sb; p.space_after = sa
    return p

def top_bar(slide):
    rect(slide, Inches(0), Inches(0), Inches(4.44), Inches(0.06), fill=RED)
    rect(slide, Inches(4.44), Inches(0), Inches(4.44), Inches(0.06), fill=BLUE)
    rect(slide, Inches(8.88), Inches(0), Inches(4.453), Inches(0.06), fill=GOLD)

def page_badge(slide, n, total=3):
    s = rect(slide, Inches(11.6), Inches(0.26), Inches(1.2), Inches(0.3), fill=BLUE)
    tf = s.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.text = f"{n:02d} / {total:02d}"
    p.font.size = Pt(8); p.font.color.rgb = WHITE; p.font.name = 'Microsoft YaHei'
    p.alignment = PP_ALIGN.CENTER

def slide_header(slide, title, sub):
    tb(slide, Inches(0.6), Inches(0.36), Inches(8), Inches(0.45), title, size=Pt(22), color=BLUE, bold=True)
    tb(slide, Inches(8.5), Inches(0.36), Inches(4), Inches(0.45), sub, size=Pt(10), color=GRAY, align=PP_ALIGN.RIGHT)
    rect(slide, Inches(0.6), Inches(0.92), Inches(12.1), Inches(0.015), fill=BLUE)

def footer(slide, left_text="乐山市沙湾区卫生健康局", right_text="2026 年 6 月"):
    rect(slide, Inches(0.6), Inches(7.05), Inches(12.1), Inches(0.008), fill=RGBColor(0xEC, 0xE9, 0xE0))
    tb(slide, Inches(0.6), Inches(7.08), Inches(5), Inches(0.22), left_text, size=Pt(7.5), color=RGBColor(0xAA,0xAA,0xAA))
    tb(slide, Inches(8.5), Inches(7.08), Inches(4), Inches(0.22), right_text, size=Pt(7.5), color=RGBColor(0xAA,0xAA,0xAA), align=PP_ALIGN.RIGHT)

def sec_label(slide, l, t, text):
    tb(slide, l, t, Inches(4), Inches(0.2), text, size=Pt(9), color=RED, bold=True)

def make_table(slide, l, t, col_w, headers, rows, font_size=Pt(10)):
    n_rows = len(rows) + 1
    n_cols = len(headers)
    total_w = sum(col_w)
    ts = slide.shapes.add_table(n_rows, n_cols, l, t, total_w, Inches(0.28 * n_rows))
    tbl = ts.table
    for ci, cw in enumerate(col_w):
        tbl.columns[ci].width = cw
    # header
    for ci, hdr in enumerate(headers):
        c = tbl.cell(0, ci); c.text = hdr
        c.fill.solid(); c.fill.fore_color.rgb = BLUE
        for p in c.text_frame.paragraphs:
            p.font.size = Pt(8.5); p.font.color.rgb = WHITE; p.font.bold = True
            p.font.name = 'Microsoft YaHei'
            p.alignment = PP_ALIGN.CENTER if ci > 0 else PP_ALIGN.LEFT
    # body
    for ri, row in enumerate(rows):
        for ci, val in enumerate(row):
            c = tbl.cell(ri+1, ci); c.text = str(val) if val is not None else "—"
            c.fill.solid()
            c.fill.fore_color.rgb = RGBColor(0xFA,0xFA,0xF8) if ri % 2 == 1 else WHITE
            for p in c.text_frame.paragraphs:
                p.font.size = font_size; p.font.color.rgb = INK
                p.font.name = 'Microsoft YaHei'
                p.alignment = PP_ALIGN.CENTER if ci > 0 else PP_ALIGN.LEFT
    return ts


# ================================================================
# SLIDE 1: 沙湾区医疗信息化现状如何？
# ================================================================
s1 = prs.slides.add_slide(blank)
top_bar(s1)
page_badge(s1, 1)
slide_header(s1, "沙湾区医疗信息化现状如何？", "乐山市沙湾区\n智慧医疗调研与规划汇报")

# -- Survey stats row --
sec_label(s1, Inches(0.6), Inches(1.12), "调研概况")
stats = [
    ("2", "公立医院（二甲）", False),
    ("5", "民营医院", False),
    ("9", "乡镇/社区机构", False),
    ("133", "村卫生室", True),
    ("14.5万", "常住人口", True),
    ("3万", "65岁以上老人", True),
]
stat_w = Inches(1.85); stat_gap = Inches(0.14); stat_y = Inches(1.38)
for i, (num, label, hl) in enumerate(stats):
    x = Inches(0.6) + i*(stat_w + stat_gap)
    rect(s1, x, stat_y, stat_w, Inches(0.72), fill=LIGHT_BG, border=BORDER)
    clr = RED if hl else BLUE
    tb(s1, x, stat_y + Inches(0.06), stat_w, Inches(0.32), num, size=Pt(20), color=clr, bold=True, align=PP_ALIGN.CENTER)
    tb(s1, x, stat_y + Inches(0.42), stat_w, Inches(0.24), label, size=Pt(9), color=GRAY, align=PP_ALIGN.CENTER)

# -- Two finding tables side by side --
find_y = Inches(2.28)

# Left table: 区级医院
sec_label(s1, Inches(0.6), find_y - Inches(0.18), "")
r1 = rect(s1, Inches(0.6), find_y, Inches(5.9), Inches(1.82), fill=WHITE, border=BORDER)
tb(s1, Inches(0.72), find_y + Inches(0.06), Inches(5.6), Inches(0.22),
   "区级医院 — 基础系统齐全，缺平台与智能化", size=Pt(11), color=BLUE, bold=True)
hosp_headers = ["机构", "现状", "缺口"]
hosp_rows = [
    ["区人民医院 [二甲]", "HIS/LIS/EMR/PACS/手麻/移动医护等齐全", "无 CDR、集成平台、BI、大数据"],
    ["区中医医院 [二级]", "新成立，同人民医院系统，新建 HIS", "其余模块共享人民医院"],
    ["区妇幼保健院 [二级]", "HIS/EMR/LIS/PACS/手麻/体检", "互联网医院、便民小程序"],
]
make_table(s1, Inches(0.72), find_y + Inches(0.34), [Inches(1.6), Inches(2.6), Inches(1.6)], hosp_headers, hosp_rows, font_size=Pt(9))

# Right table: 基层机构
sec_label(s1, Inches(6.8), find_y - Inches(0.18), "")
r2 = rect(s1, Inches(6.8), find_y, Inches(5.9), Inches(1.82), fill=WHITE, border=BORDER)
tb(s1, Inches(6.92), find_y + Inches(0.06), Inches(5.6), Inches(0.22),
   "基层机构 — 无自建系统，信息孤岛严重", size=Pt(11), color=BLUE, bold=True)
grass_headers = ["维度", "现状与问题"]
grass_rows = [
    ["系统情况", "除福禄镇外均无自建系统，依赖省统建云平台"],
    ["村卫生室", "133 个村卫生室无任何信息系统"],
    ["数据互通", "基层系统与公卫云数据不通，重复录入"],
    ["转诊量", "年转诊约 100+ 人（福禄镇 200+），双向转诊无反馈"],
    ["信息化投入", "仅福禄镇自建系统，年费约 10 万元"],
]
make_table(s1, Inches(6.92), find_y + Inches(0.34), [Inches(1.2), Inches(4.5)], grass_headers, grass_rows, font_size=Pt(9))

# -- Four pain points --
pain_y = Inches(4.25)
sec_label(s1, Inches(0.6), pain_y, "核心痛点")
pains = [
    ("01", "数据孤岛", "区级缺CDR与集成平台，\n基层与公卫数据不通", RED, TAG_RED),
    ("02", "协同断层", "双向转诊无反馈，远程会诊\n未建立，检查结果不互认", BLUE, TAG_BLUE),
    ("03", "基层薄弱", "133个村卫生室零信息化，\n8个乡镇无自建系统", GOLD, TAG_GOLD),
    ("04", "智慧缺失", "无AI辅助、无CDSS、\n无BI决策，距智慧三星差距大", RED, TAG_RED),
]
pain_w = Inches(2.9); pain_gap = Inches(0.13); pain_y2 = Inches(4.5)
for i, (num, name, desc, accent, tag_c) in enumerate(pains):
    x = Inches(0.6) + i*(pain_w + pain_gap)
    rect(s1, x, pain_y2, pain_w, Inches(1.85), fill=WHITE, border=BORDER)
    c = oval(s1, x + Inches(1.1), pain_y2 + Inches(0.12), Inches(0.5), Inches(0.5), fill=accent)
    tf = c.text_frame; p = tf.paragraphs[0]
    p.text = num; p.font.size = Pt(14); p.font.color.rgb = WHITE; p.font.bold = True
    p.font.name = 'Microsoft YaHei'; p.alignment = PP_ALIGN.CENTER
    tb(s1, x + Inches(0.1), pain_y2 + Inches(0.72), Inches(2.7), Inches(0.22), name, size=Pt(12), color=INK, bold=True, align=PP_ALIGN.CENTER)
    tb(s1, x + Inches(0.15), pain_y2 + Inches(1.0), Inches(2.6), Inches(0.48), desc, size=Pt(9), color=DARK_GRAY, align=PP_ALIGN.CENTER)

footer(s1)


# ================================================================
# SLIDE 2: 怎么建智慧医疗？
# ================================================================
s2 = prs.slides.add_slide(blank)
top_bar(s2)
page_badge(s2, 2)
slide_header(s2, "怎么建智慧医疗？", "乐山市沙湾区\n智慧医疗调研与规划汇报")

sec_label(s2, Inches(0.6), Inches(1.12), "建设策略")
# Legend row
legend_y = Inches(1.38)
for i, (label, dot_color, x_off) in enumerate([
    ("一期统建", RED, 0), ("二期建设", BLUE, 1.8), ("已建/利旧对接", GRAY, 3.6)
]):
    oval(s2, Inches(0.7 + x_off), legend_y + Inches(0.02), Inches(0.15), Inches(0.15), fill=dot_color)
    tb(s2, Inches(0.9 + x_off), legend_y, Inches(1.2), Inches(0.22), label, size=Pt(9), color=INK)
tb(s2, Inches(8.5), legend_y, Inches(4), Inches(0.22), "总投资量级：5,690 万元", size=Pt(10), color=INK, bold=True, align=PP_ALIGN.RIGHT)

# ---- Two-column module layout with ALIGNED grid ----
# Use a unified row-height grid. Each "row" in the layout has the same Y across columns.
col_l = Inches(0.6)
col_r = Inches(6.8)
col_w_half = Inches(5.9)

# Row grid (top Y of each module group block):
# Row 0: 区域大数据平台 (left)  |  区级医院能力提升 (right)   -- y=1.75
# Row 1: 基层服务能力提升 (left)  |  医院集成平台+大数据 (right)  -- y=3.85
# Row 2: 区域协同应用 (left)      |  综合决策监管 (right)       -- y=4.95
# Row 3:                           |  急诊急救协同 (right)       -- y=5.75

def module_block(slide, x, y, w, title, items, head_color=RED):
    """Draw a module group block. items is list of (bold_label, [tag_items]).
       Returns the bottom Y of this block."""
    n = len(items)
    h = Inches(0.3) + Inches(0.28 * n) + Inches(0.08)
    rect(slide, x, y, w, h, fill=WHITE, border=BORDER)
    # header
    rect(slide, x, y, w, Inches(0.28), fill=head_color)
    tb(slide, x + Inches(0.08), y + Inches(0.02), w - Inches(0.16), Inches(0.24),
       title, size=Pt(9), color=WHITE, bold=True, align=PP_ALIGN.LEFT)
    # body
    body_y = y + Inches(0.32)
    for j, (label, tags) in enumerate(items):
        row_y = body_y + j * Inches(0.26)
        tf = rich_tb(slide, x + Inches(0.08), row_y, w - Inches(0.16), Inches(0.24))
        p = tf.paragraphs[0]
        run(p, label, size=Pt(9), color=INK, bold=True)
        run(p, "  " + " · ".join(tags), size=Pt(8.5), color=DARK_GRAY)
    return y + h

# Unified grid Y positions for left & right columns
GRID = [
    Inches(1.72),  # row 0
    Inches(3.78),  # row 1
    Inches(4.88),  # row 2
    Inches(5.72),  # row 3
]

# Left column blocks
module_block(s2, col_l, GRID[0], col_w_half, "区域医疗大数据平台（一期统建）", [
    ("基础平台:", ["主索引与注册", "统一认证与权限", "总线服务", "API网关", "信息标准管理"]),
    ("数据服务:", ["全方式采集", "加密脱敏", "数据治理", "质量管理", "完整性管理"]),
    ("资源中心:", ["基础资源库", "医疗卫生资源库", "电子病历库(58数据集)", "健康档案库"]),
    ("集成门户:", ["统一门户", "单点登录", "机构/人员/设备管理", "互联互通共享"]),
], head_color=RED)

module_block(s2, col_l, GRID[1], col_w_half, "基层服务能力提升", [
    ("基本医疗:", ["全结构化电子病历(一期)", "区域PACS+AI肺结节(一期)", "区域LIS(二期)"]),
    ("公卫服务:", ["妇幼/计免/公卫(已建对接)", "健康档案统计(一期)", "村医随访系统(二期)"]),
    ("质量管理:", ["AI病历质控(一期)", "基层CDSS临床辅助决策(一期)"]),
], head_color=BLUE)

module_block(s2, col_l, GRID[2], col_w_half, "区域协同应用", [
    ("便民协同:", ["区域互联网医院(一期)", "电子健康卡(一期)", "统一预约/支付/查询(一期)"]),
    ("医疗协同:", ["区域双向转诊(一期)", "处方前置审核(二期)", "区域DRG(二期)", "检查互认(市建利旧)"]),
    ("管理协同:", ["区域心电/影像(利旧)", "消毒供应中心(一期)", "HRP/物资/采购(二期)", "ChatBI AI卫健监管(二期)"]),
], head_color=GOLD)

# Right column blocks
module_block(s2, col_r, GRID[0], col_w_half, "区级医院能力提升（电子病历5级目标）", [
    ("临床核心(二期):", ["门急诊/住院工作站", "移动医护", "临床路径", "患者全息视图", "业务闭环"]),
    ("医疗管理(一/二期):", ["智慧医务", "院感/不良事件", "护理/输血/重症", "单病种/全病程", "CA认证"]),
    ("智慧病房(二期):", ["床旁交互", "移动护理", "物联网管控", "护士站调度"]),
    ("智慧中医(二期):", ["智能辅助诊疗", "治未病平台"]),
], head_color=RED)

module_block(s2, col_r, GRID[1], col_w_half, "医院集成平台 + 大数据中心（一期）", [
    ("平台与数据:", ["医院互联互通平台", "CDR临床数据中心", "国考绩效指标管理", "等级医院评审指标"]),
], head_color=BLUE)

module_block(s2, col_r, GRID[2], col_w_half, "综合决策监管（一期统建）", [
    ("监管看板:", ["运营监控(门诊/住院人次、均费)", "服务效能(医保/收支)", "协同监控(转诊/影像/检验)", "数据资源监控(采集/治理)"]),
], head_color=GOLD)

module_block(s2, col_r, GRID[3], col_w_half, "急诊急救协同（一期）", [
    ("急救中心:", ["胸痛中心", "卒中中心", "创伤中心", "智慧急救平台", "孕产/新生儿中心(二期)"]),
], head_color=RED)

# -- Goal bar --
goal_y = Inches(6.2)
goal_shape = rect(s2, Inches(0.6), goal_y, Inches(12.1), Inches(0.55), fill=BLUE)
tf = goal_shape.text_frame; tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "目标：电子病历 5 级 + 互联互通四级甲等 + 智慧三星 —— 2028 年交付"
p.font.size = Pt(14); p.font.color.rgb = WHITE; p.font.bold = True
p.font.name = 'Microsoft YaHei'; p.alignment = PP_ALIGN.CENTER

footer(s2)


# ================================================================
# SLIDE 3: 花多少钱？值不值？
# ================================================================
s3 = prs.slides.add_slide(blank)
top_bar(s3)
page_badge(s3, 3)
slide_header(s3, "花多少钱？值不值？", "乐山市沙湾区\n智慧医疗调研与规划汇报")

# Left: Investment table
sec_label(s3, Inches(0.6), Inches(1.12), "投资概览")
tb(s3, Inches(0.6), Inches(1.32), Inches(7), Inches(0.28),
   "智慧医疗总投资 5,690 万元", size=Pt(15), color=INK, bold=True)

inv_headers = ["建设内容", "方式", "分期", "预算"]
inv_rows = [
    ["区域医疗大数据平台",        "统建",     "一期",    "含在总额"],
    ["智慧医疗综合决策监管",      "统建",     "一期",    "含在总额"],
    ["基层服务能力提升",          "统建",     "一/二期", "含在总额"],
    ["区级医院能力提升(含5级)",   "统建",     "一/二期", "含在总额"],
    ["医院集成平台+大数据中心",   "统建",     "一期",    "含在总额"],
    ["区域协同应用",              "统建+利旧","一/二期", "含在总额"],
    ["系统接口与对接",            "统建",     "一期",    "含在总额"],
]
inv_col_w = [Inches(3.0), Inches(1.2), Inches(1.0), Inches(1.1)]
ts = make_table(s3, Inches(0.6), Inches(1.72), inv_col_w, inv_headers, inv_rows, font_size=Pt(9))

# Sum row
sum_y = Inches(1.72) + Inches(0.28 * 8)
sum_w = sum(inv_col_w)
s_sum = rect(s3, Inches(0.6), sum_y, sum_w, Inches(0.28), fill=SUM_BG, border=BLUE, border_w=Pt(1.5))
tb(s3, Inches(0.7),  sum_y + Inches(0.02), Inches(2.8), Inches(0.24), "合计", size=Pt(9), color=INK, bold=True)
tb(s3, Inches(3.7),  sum_y + Inches(0.02), Inches(1.1), Inches(0.24), "", size=Pt(9))
tb(s3, Inches(4.85), sum_y + Inches(0.02), Inches(0.9), Inches(0.24), "", size=Pt(9))
tb(s3, Inches(5.8),  sum_y + Inches(0.02), Inches(0.9), Inches(0.24), "5,690 万元", size=Pt(9), color=INK, bold=True, align=PP_ALIGN.RIGHT)

# Note
note_y = sum_y + Inches(0.4)
note_s = rect(s3, Inches(0.6), note_y, Inches(6.15), Inches(0.4), fill=NOTE_BG, border=RGBColor(0xE5,0xD9,0xA0))
tf = note_s.text_frame; tf.word_wrap = True
p = tf.paragraphs[0]
run(p, "资金说明：", size=Pt(8.5), color=GOLD, bold=True)
run(p, "统建部分做预算；已建/自建系统利旧对接不另计费。区人民医院新院区专项预算 2,400 万元（2028 年交付）。", size=Pt(8.5), color=GOLD)

# Right: Outcomes
out_x = Inches(7.2)
sec_label(s3, out_x, Inches(1.12), "预期成效")
out_box = rect(s3, out_x, Inches(1.32), Inches(5.5), Inches(5.2), fill=WHITE, border=BORDER)
tb(s3, out_x + Inches(0.15), Inches(1.42), Inches(5), Inches(0.25),
   "建设前后对比", size=Pt(13), color=BLUE, bold=True)

outcomes = [
    ("1", "数据贯通，一网共享",
     "建成区域大数据平台+CDR，终结\"烟囱\"时代，全区医疗机构病历、检验、影像、档案互联互通。"),
    ("2", "基层从\"零\"到\"有\"",
     "133个村卫生室+8个乡镇首获电子病历、PACS、CDSS和AI辅助，基层诊疗能力质的飞跃。"),
    ("3", "区级达\"5-4-3\"标准",
     "电子病历5级+互联互通四甲+智慧三星，跻身全省区县级医院第一梯队。"),
    ("4", "便民惠民全覆盖",
     "互联网医院、电子健康卡、统一预约/支付/查询，居民\"一部手机管健康\"。"),
    ("5", "AI赋能精准决策",
     "AI病历质控+CDSS辅助+ChatBI卫健监管，从\"经验驱动\"到\"数据驱动\"。"),
]
for j, (num, title, desc) in enumerate(outcomes):
    yy = Inches(1.85) + j * Inches(0.88)
    # number circle
    c = oval(s3, out_x + Inches(0.2), yy, Inches(0.3), Inches(0.3), fill=RED)
    tf = c.text_frame; p = tf.paragraphs[0]
    p.text = num; p.font.size = Pt(10); p.font.color.rgb = WHITE; p.font.bold = True
    p.font.name = 'Microsoft YaHei'; p.alignment = PP_ALIGN.CENTER
    # text
    tb(s3, out_x + Inches(0.62), yy - Inches(0.02), Inches(4.6), Inches(0.18), title, size=Pt(10), color=INK, bold=True)
    tb(s3, out_x + Inches(0.62), yy + Inches(0.18), Inches(4.6), Inches(0.5), desc, size=Pt(8.5), color=DARK_GRAY)
    # separator
    if j < 4:
        rect(s3, out_x + Inches(0.62), yy + Inches(0.72), Inches(4.8), Inches(0.005), fill=RGBColor(0xEE,0xEE,0xEE))

footer(s3)


# === Save ===
out = r"E:\claude\output\沙湾智慧医疗调研规划_3页.pptx"
prs.save(out)
print(f"Saved: {out}")
